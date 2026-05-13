"""PowerPoint presentation generator for earnings season summaries."""

import os
from datetime import date
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .config import Company


# Brand colors matching the email template
DARK_BLUE = RGBColor(0x1A, 0x36, 0x5D)    # #1a365d - header
ACCENT_BLUE = RGBColor(0x31, 0x82, 0xCE)   # #3182ce - accents
LIGHT_BLUE = RGBColor(0xEB, 0xF8, 0xFF)    # #ebf8ff - light bg
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x2D, 0x37, 0x48)     # #2d3748
MED_GRAY = RGBColor(0x4A, 0x55, 0x68)      # #4a5568
LIGHT_GRAY = RGBColor(0xF7, 0xFA, 0xFC)    # #f7fafc
GREEN = RGBColor(0x38, 0xA1, 0x69)          # #38a169
RED = RGBColor(0xE5, 0x3E, 0x3E)            # #e53e3e
PURPLE = RGBColor(0x7C, 0x3A, 0xED)         # #7c3aed - PE activity


class SeasonalPPTGenerator:
    """Generates branded PowerPoint presentations for earnings season summaries."""

    def __init__(self):
        self.prs = Presentation()
        # Set slide dimensions to widescreen 16:9
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def generate(
        self,
        quarter: str,
        season_data: dict,
        output_path: str
    ) -> str:
        """Generate the earnings season summary presentation.

        Args:
            quarter: e.g., "Q1 2025"
            season_data: Dict containing all season data with keys:
                - first_report_date, last_report_date
                - executive_summary (str)
                - sector_themes (str)
                - company_highlights (dict: company_name -> str)
                - outlook (str)
                - market_data (list of dicts with company stock data)
                - hyperscaler_summary (str)
                - pe_summary (str)
                - companies_reported (list of company names)
            output_path: Path to save the .pptx file.

        Returns:
            The output_path where the file was saved.
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

        first_date = season_data.get("first_report_date")
        last_date = season_data.get("last_report_date")

        self._add_title_slide(quarter, first_date, last_date)
        self._add_executive_summary_slide(
            quarter, season_data.get("executive_summary", "")
        )
        self._add_market_performance_slide(
            quarter, season_data.get("market_data", [])
        )
        self._add_sector_themes_slide(
            quarter, season_data.get("sector_themes", "")
        )

        # Per-company slides
        company_highlights = season_data.get("company_highlights", {})
        for company_name, highlights in company_highlights.items():
            self._add_company_slide(company_name, highlights)

        self._add_hyperscaler_slide(
            quarter, season_data.get("hyperscaler_summary", "")
        )
        self._add_pe_activity_slide(
            quarter, season_data.get("pe_summary", "")
        )
        self._add_outlook_slide(
            quarter, season_data.get("outlook", "")
        )

        # Save presentation
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        self.prs.save(output_path)
        return output_path

    def _add_dark_background(self, slide):
        """Add a dark blue background to a slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BLUE

    def _add_light_background(self, slide):
        """Add a light background to a slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

    def _add_accent_bar(self, slide, top=Inches(0), width=None):
        """Add a thin accent bar at the top of a slide."""
        if width is None:
            width = self.prs.slide_width
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), top, width, Inches(0.08)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_BLUE
        shape.line.fill.background()

    def _add_header_bar(self, slide, title_text):
        """Add a consistent header bar with title to content slides."""
        # Header background
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(1.1)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = DARK_BLUE
        header.line.fill.background()

        # Title text in header
        txBox = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.2),
            Inches(12), Inches(0.7)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE
        p.font.bold = True

        # Accent bar below header
        self._add_accent_bar(slide, top=Inches(1.1))

    def _add_bullet_content(self, text_frame, content_text, font_size=Pt(14)):
        """Add bullet-point content from a text string to a text frame."""
        text_frame.word_wrap = True
        lines = content_text.strip().split("\n")

        first = True
        for line in lines:
            line = line.strip()
            if not line:
                continue

            if first:
                p = text_frame.paragraphs[0]
                first = False
            else:
                p = text_frame.add_paragraph()

            # Handle markdown-style bullets
            if line.startswith("- ") or line.startswith("* "):
                p.text = line[2:]
                p.level = 0
                p.space_before = Pt(4)
            elif line.startswith("  - ") or line.startswith("  * "):
                p.text = line[4:]
                p.level = 1
                p.space_before = Pt(2)
            elif line.startswith("**") and line.endswith("**"):
                # Bold section header
                p.text = line.strip("*")
                p.font.bold = True
                p.space_before = Pt(10)
            elif line.startswith("### "):
                p.text = line[4:]
                p.font.bold = True
                p.font.size = Pt(16)
                p.space_before = Pt(12)
            elif line.startswith("## "):
                p.text = line[3:]
                p.font.bold = True
                p.font.size = Pt(18)
                p.space_before = Pt(14)
            else:
                p.text = line
                p.space_before = Pt(4)

            p.font.size = font_size
            p.font.color.rgb = DARK_GRAY

    def _add_title_slide(
        self,
        quarter: str,
        first_date: Optional[date],
        last_date: Optional[date]
    ):
        """Add the title slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        self._add_dark_background(slide)

        # Accent bar
        self._add_accent_bar(slide, top=Inches(2.8))

        # Main title
        txBox = slide.shapes.add_textbox(
            Inches(1), Inches(1.5),
            Inches(11.333), Inches(1.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{quarter} Earnings Season Summary"
        p.font.size = Pt(44)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Subtitle with date range
        txBox2 = slide.shapes.add_textbox(
            Inches(1), Inches(3.2),
            Inches(11.333), Inches(1.0)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]

        if first_date and last_date:
            date_range = (
                f"{first_date.strftime('%B %d')} - "
                f"{last_date.strftime('%B %d, %Y')}"
            )
            p2.text = f"Reporting Period: {date_range}"
        else:
            p2.text = "Data Center Infrastructure Companies"

        p2.font.size = Pt(20)
        p2.font.color.rgb = ACCENT_BLUE
        p2.alignment = PP_ALIGN.CENTER

        # Subtitle line 2
        txBox3 = slide.shapes.add_textbox(
            Inches(1), Inches(4.2),
            Inches(11.333), Inches(0.6)
        )
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = "Data Center Infrastructure Supply Chain Intelligence"
        p3.font.size = Pt(16)
        p3.font.color.rgb = RGBColor(0x71, 0x80, 0x96)  # Muted gray
        p3.alignment = PP_ALIGN.CENTER

        # Footer
        txBox4 = slide.shapes.add_textbox(
            Inches(1), Inches(6.5),
            Inches(11.333), Inches(0.5)
        )
        tf4 = txBox4.text_frame
        p4 = tf4.paragraphs[0]
        p4.text = "Generated by Company Tracker"
        p4.font.size = Pt(11)
        p4.font.color.rgb = RGBColor(0x71, 0x80, 0x96)
        p4.alignment = PP_ALIGN.CENTER

    def _add_executive_summary_slide(self, quarter: str, summary_text: str):
        """Add the executive summary slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._add_light_background(slide)
        self._add_header_bar(slide, f"{quarter} Executive Summary")

        # Content area
        txBox = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.4),
            Inches(12.1), Inches(5.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        self._add_bullet_content(tf, summary_text)

    def _add_market_performance_slide(self, quarter: str, market_data: list):
        """Add market performance table slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._add_light_background(slide)
        self._add_header_bar(slide, f"{quarter} Market Performance")

        if not market_data:
            txBox = slide.shapes.add_textbox(
                Inches(1), Inches(2.5),
                Inches(11), Inches(1)
            )
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = "No market performance data available for this period."
            p.font.size = Pt(16)
            p.font.color.rgb = MED_GRAY
            p.alignment = PP_ALIGN.CENTER
            return

        # Create table
        rows = len(market_data) + 1  # +1 for header
        cols = 5
        table_width = Inches(11.5)
        table_height = Inches(min(0.35 * rows + 0.1, 5.8))

        table_shape = slide.shapes.add_table(
            rows, cols,
            Inches(0.9), Inches(1.4),
            table_width, table_height
        )
        table = table_shape.table

        # Set column widths
        col_widths = [Inches(3.0), Inches(1.5), Inches(2.5), Inches(2.5), Inches(2.0)]
        for i, width in enumerate(col_widths):
            table.columns[i].width = width

        # Header row
        headers = ["Company", "Ticker", "Season Start Price", "Season End Price", "Change %"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Data rows
        for row_idx, data in enumerate(market_data, start=1):
            bg_color = WHITE if row_idx % 2 == 1 else LIGHT_GRAY

            company_cell = table.cell(row_idx, 0)
            company_cell.text = data.get("company", "")
            company_cell.fill.solid()
            company_cell.fill.fore_color.rgb = bg_color

            ticker_cell = table.cell(row_idx, 1)
            ticker_cell.text = data.get("ticker", "N/A")
            ticker_cell.fill.solid()
            ticker_cell.fill.fore_color.rgb = bg_color

            start_price = data.get("start_price")
            start_cell = table.cell(row_idx, 2)
            start_cell.text = f"${start_price:.2f}" if start_price else "N/A"
            start_cell.fill.solid()
            start_cell.fill.fore_color.rgb = bg_color

            end_price = data.get("end_price")
            end_cell = table.cell(row_idx, 3)
            end_cell.text = f"${end_price:.2f}" if end_price else "N/A"
            end_cell.fill.solid()
            end_cell.fill.fore_color.rgb = bg_color

            change = data.get("change_percent")
            change_cell = table.cell(row_idx, 4)
            if change is not None:
                sign = "+" if change >= 0 else ""
                change_cell.text = f"{sign}{change:.1f}%"
                p = change_cell.text_frame.paragraphs[0]
                p.font.color.rgb = GREEN if change >= 0 else RED
                p.font.bold = True
            else:
                change_cell.text = "N/A"
            change_cell.fill.solid()
            change_cell.fill.fore_color.rgb = bg_color

            # Format all cells in row
            for col_idx in range(cols):
                cell = table.cell(row_idx, col_idx)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(10)
                if col_idx > 0:
                    p.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _add_sector_themes_slide(self, quarter: str, themes_text: str):
        """Add sector themes slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._add_light_background(slide)
        self._add_header_bar(slide, f"{quarter} Sector Themes")

        txBox = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.4),
            Inches(12.1), Inches(5.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        self._add_bullet_content(tf, themes_text)

    def _add_company_slide(self, company_name: str, highlights_text: str):
        """Add a per-company detail slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._add_light_background(slide)
        self._add_header_bar(slide, company_name)

        txBox = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.4),
            Inches(12.1), Inches(5.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        self._add_bullet_content(tf, highlights_text)

    def _add_hyperscaler_slide(self, quarter: str, summary_text: str):
        """Add hyperscaler activity summary slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._add_light_background(slide)
        self._add_header_bar(slide, f"{quarter} Hyperscaler Activity")

        # Blue accent sidebar
        sidebar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(1.18),
            Inches(0.08), Inches(6.32)
        )
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = RGBColor(0x0E, 0xA5, 0xE9)  # Sky blue
        sidebar.line.fill.background()

        txBox = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.4),
            Inches(12.1), Inches(5.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        self._add_bullet_content(tf, summary_text)

    def _add_pe_activity_slide(self, quarter: str, summary_text: str):
        """Add PE activity summary slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._add_light_background(slide)
        self._add_header_bar(slide, f"{quarter} Private Equity Activity")

        # Purple accent sidebar
        sidebar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(1.18),
            Inches(0.08), Inches(6.32)
        )
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = PURPLE
        sidebar.line.fill.background()

        txBox = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.4),
            Inches(12.1), Inches(5.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        self._add_bullet_content(tf, summary_text)

    def _add_outlook_slide(self, quarter: str, outlook_text: str):
        """Add the outlook/forward guidance slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._add_light_background(slide)
        self._add_header_bar(slide, f"{quarter} Outlook & Forward Guidance")

        txBox = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.4),
            Inches(12.1), Inches(5.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        self._add_bullet_content(tf, outlook_text)
